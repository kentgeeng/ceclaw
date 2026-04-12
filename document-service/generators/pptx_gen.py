import httpx
import json
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

from config import LLM_URL, LLM_MODEL

SEARXNG_URL = "http://127.0.0.1:2337/v1/search"

# ── 色票 ──────────────────────────────────────────
C = {
    "dark": {
        "bg":     RGBColor(0x0A,0x0A,0x14), "sidebar": RGBColor(0x14,0x10,0x2E),
        "acc":    RGBColor(0xE8,0x40,0x5A), "acc2":    RGBColor(0x6C,0x63,0xFF),
        "title":  RGBColor(0xFF,0xFF,0xFF), "body":    RGBColor(0xCC,0xCC,0xDD),
        "sub":    RGBColor(0x88,0x88,0xAA), "rule":    RGBColor(0x2A,0x2A,0x44),
        "num_fg": RGBColor(0xFF,0xFF,0xFF), "white":   RGBColor(0xFF,0xFF,0xFF),
    },
    "light": {
        "bg":     RGBColor(0xF6,0xF6,0xFC), "sidebar": RGBColor(0x1A,0x56,0xDB),
        "acc":    RGBColor(0x1A,0x56,0xDB), "acc2":    RGBColor(0x7C,0x3A,0xED),
        "title":  RGBColor(0x0D,0x0D,0x2E), "body":    RGBColor(0x22,0x22,0x44),
        "sub":    RGBColor(0x66,0x66,0x88), "rule":    RGBColor(0xD8,0xD8,0xEE),
        "num_fg": RGBColor(0xFF,0xFF,0xFF), "white":   RGBColor(0xFF,0xFF,0xFF),
    },
    "taiwan": {
        "bg":     RGBColor(0x0E,0x0E,0x12), "sidebar": RGBColor(0xBF,0x0A,0x1E),
        "acc":    RGBColor(0xE8,0x15,0x2A), "acc2":    RGBColor(0xFF,0xA5,0x00),
        "title":  RGBColor(0xFF,0xFF,0xFF), "body":    RGBColor(0xE8,0xE0,0xD0),
        "sub":    RGBColor(0xAA,0x99,0x88), "rule":    RGBColor(0x2A,0x1A,0x1A),
        "num_fg": RGBColor(0xFF,0xFF,0xFF), "white":   RGBColor(0xFF,0xFF,0xFF),
    },
}


async def search_web(query: str, num: int = 6) -> list:
    try:
        async with httpx.AsyncClient(timeout=15) as c:
            r = await c.post(SEARXNG_URL, json={"query": query})
            results = r.json().get("data", {}).get("web", [])[:num]
            return [{"title": x["title"], "description": x.get("description", "")} for x in results]
    except Exception:
        return []


async def llm_generate_slides(topic: str, pages: int, search_results: list) -> list:
    context = "\n".join([f"- {r['title']}: {r['description'][:200]}" for r in search_results])
    prompt = f"""你是資深商業顧問，製作高品質繁體中文簡報。

主題：{topic}
頁數：{pages} 頁（含封面）

網路最新參考資料：
{context}

規則：
- 每頁標題 10 字以內，簡潔有力
- 每頁 3-4 要點，每點 20-40 字，含具體數據或案例
- 第一頁封面只要 title 和 subtitle
- 台灣本地化視角，禁止空話

只輸出 JSON 陣列：
[
  {{"type":"cover","title":"主標題","subtitle":"副標題"}},
  {{"type":"content","title":"頁面標題","points":["要點1","要點2","要點3"]}}
]"""
    async with httpx.AsyncClient(timeout=90) as c:
        r = await c.post(LLM_URL, json={
            "model": LLM_MODEL,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0, "max_tokens": 3000,
        })
    text = r.json()["choices"][0]["message"]["content"].strip()
    text = re.sub(r'<\|channel\|>.*?<channel\|>', '', text, flags=re.DOTALL)
    text = re.sub(r'<think>.*?</think>', '', text, flags=re.DOTALL)
    return json.loads(text[text.find("["):text.rfind("]")+1])


# ── 工具 ──────────────────────────────────────────
def bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def rect(slide, x, y, w, h, color, alpha=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, font="Noto Serif CJK TC"):
    t = slide.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.name = font
    return t


# ══════════════════════════════════════════════════
# THEME: DARK PRO
# ══════════════════════════════════════════════════
def dark_cover(prs, data, t):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    bg(slide, t["bg"])

    # 左側漸層色塊（深紫）
    rect(slide, 0, 0, Inches(4.2), H, t["sidebar"])
    # 左側底部紅色細條
    rect(slide, 0, H - Inches(0.12), Inches(4.2), Inches(0.12), t["acc"])
    # 右上角裝飾三角區塊（用兩個矩形模擬）
    rect(slide, W - Inches(1.5), 0, Inches(1.5), Inches(1.5), t["acc"])
    rect(slide, W - Inches(0.8), 0, Inches(0.8), Inches(0.8), t["bg"])

    # CECLAW logo
    txt(slide, "CECLAW", Inches(0.3), Inches(0.28), Inches(3.6), Inches(0.5),
        12, t["white"], bold=True)
    txt(slide, "Enterprise AI Platform", Inches(0.3), Inches(0.7), Inches(3.6), Inches(0.4),
        8, RGBColor(0xAA,0xAA,0xCC))

    # 左側大數字裝飾
    txt(slide, "01", Inches(0.2), H - Inches(1.4), Inches(2), Inches(1.0),
        52, RGBColor(0xFF,0xFF,0xFF,), bold=True)

    # 主標題（右側）
    txt(slide, data.get("title",""), Inches(4.6), Inches(1.5), Inches(5.0), Inches(2.8),
        30, t["title"], bold=True)

    # 紅色分隔線
    rect(slide, Inches(4.6), Inches(4.5), Inches(2.0), Inches(0.05), t["acc"])

    # 副標題
    txt(slide, data.get("subtitle",""), Inches(4.6), Inches(4.7), Inches(5.0), Inches(0.8),
        13, t["sub"])

    # 右下 ColdElectric
    txt(slide, "ColdElectric · 數據主權 · 本地推理", W - Inches(5.5), H - Inches(0.45),
        Inches(5.3), Inches(0.35), 9, t["sub"], align=PP_ALIGN.RIGHT)


def dark_content(prs, data, t, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    bg(slide, t["bg"])

    # 頂部色條
    rect(slide, 0, 0, W, Inches(0.055), t["acc"])
    # 頂部右側色塊
    rect(slide, W - Inches(2.5), 0, Inches(2.5), Inches(0.8), t["rule"])

    # 標題區
    rect(slide, Inches(0.45), Inches(0.22), Inches(0.05), Inches(0.55), t["acc"])
    txt(slide, data.get("title",""), Inches(0.62), Inches(0.18), Inches(7.0), Inches(0.72),
        24, t["title"], bold=True)

    # 頁碼（右上）
    txt(slide, f"{page_num:02d}/{total:02d}", W - Inches(2.3), Inches(0.22),
        Inches(1.8), Inches(0.55), 11, t["sub"], align=PP_ALIGN.RIGHT)

    # 橫線
    rect(slide, Inches(0.45), Inches(0.95), Inches(9.1), Inches(0.018), t["rule"])

    # 要點
    points = data.get("points", [])
    colors = [t["acc"], t["acc2"], t["acc"], t["acc2"]]
    y_base = Inches(1.15)
    gap = (H - Inches(1.8)) / max(len(points[:4]), 1)

    for i, pt in enumerate(points[:4]):
        y = y_base + i * gap
        # 左側色條
        rect(slide, Inches(0.45), y + Inches(0.08), Inches(0.06), gap - Inches(0.18), colors[i])
        # 數字
        txt(slide, str(i+1), Inches(0.58), y + Inches(0.06), Inches(0.32), Inches(0.32),
            11, colors[i], bold=True)
        # 要點文字
        txt(slide, pt, Inches(0.98), y, Inches(8.5), gap - Inches(0.08),
            13, t["body"])

    # 底部
    rect(slide, 0, H - Inches(0.08), W, Inches(0.08), t["rule"])
    txt(slide, "CECLAW · ColdElectric", Inches(0.45), H - Inches(0.38),
        Inches(4), Inches(0.3), 8, t["sub"])


# ══════════════════════════════════════════════════
# THEME: WHITE CLEAN
# ══════════════════════════════════════════════════
def light_cover(prs, data, t):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    bg(slide, t["bg"])

    # 頂部藍色寬條
    rect(slide, 0, 0, W, Inches(2.2), t["sidebar"])
    # 頂部條底部裝飾線
    rect(slide, 0, Inches(2.2), W, Inches(0.04), t["acc2"])

    # CECLAW（白字在藍條上）
    txt(slide, "CECLAW", Inches(0.6), Inches(0.35), Inches(4), Inches(0.55),
        13, RGBColor(0xFF,0xFF,0xFF), bold=True)
    txt(slide, "Enterprise AI Platform", Inches(0.6), Inches(0.82), Inches(4), Inches(0.4),
        9, RGBColor(0xBB,0xCC,0xFF))

    # 主標題（藍條下方）
    txt(slide, data.get("title",""), Inches(0.6), Inches(2.55), Inches(8.8), Inches(2.5),
        32, t["title"], bold=True)

    # 藍色短線
    rect(slide, Inches(0.6), Inches(5.2), Inches(1.8), Inches(0.055), t["acc"])

    # 副標題
    txt(slide, data.get("subtitle",""), Inches(0.6), Inches(5.4), Inches(8.8), Inches(0.6),
        13, t["sub"])

    # 底部品牌列
    rect(slide, 0, H - Inches(0.5), W, Inches(0.5), t["rule"])
    txt(slide, "ColdElectric · 數據主權 · 本地推理", Inches(0.6), H - Inches(0.42),
        Inches(8), Inches(0.35), 9, t["sub"])


def light_content(prs, data, t, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    bg(slide, t["bg"])

    # 頂部藍條（細）
    rect(slide, 0, 0, W, Inches(0.08), t["sidebar"])

    # 標題背景淡色
    rect(slide, 0, Inches(0.08), W, Inches(0.88), RGBColor(0xE8,0xEE,0xFF))
    txt(slide, data.get("title",""), Inches(0.5), Inches(0.12), Inches(7.5), Inches(0.78),
        22, t["title"], bold=True)
    txt(slide, f"{page_num:02d} / {total:02d}", W - Inches(1.8), Inches(0.22),
        Inches(1.6), Inches(0.55), 11, t["sub"], align=PP_ALIGN.RIGHT)

    # 橫線
    rect(slide, Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.025), t["rule"])

    # 要點（圓角卡片風格，用矩形模擬）
    points = data.get("points", [])
    colors = [t["acc"], t["acc2"], t["acc"], t["acc2"]]
    y_base = Inches(1.15)
    card_h = Inches(1.2)
    gap = Inches(0.12)

    for i, pt in enumerate(points[:4]):
        y = y_base + i * (card_h + gap)
        # 卡片背景
        rect(slide, Inches(0.5), y, Inches(9.0), card_h, RGBColor(0xFF,0xFF,0xFF))
        # 左側色條
        rect(slide, Inches(0.5), y, Inches(0.08), card_h, colors[i])
        # 數字
        txt(slide, str(i+1), Inches(0.68), y + Inches(0.1), Inches(0.36), Inches(0.36),
            12, colors[i], bold=True)
        # 要點文字
        txt(slide, pt, Inches(1.12), y + Inches(0.08), Inches(8.2), card_h - Inches(0.16),
            12, t["body"])

    # 底部
    rect(slide, 0, H - Inches(0.42), W, Inches(0.42), RGBColor(0xE8,0xEE,0xFF))
    txt(slide, "CECLAW · ColdElectric", Inches(0.5), H - Inches(0.35),
        Inches(4), Inches(0.28), 8, t["sub"])


# ══════════════════════════════════════════════════
# THEME: TAIWAN RED
# ══════════════════════════════════════════════════
def taiwan_cover(prs, data, t):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    bg(slide, t["bg"])

    # 右側大紅色塊（對角設計）
    rect(slide, Inches(5.8), 0, Inches(4.2), H, t["sidebar"])
    # 紅色塊左邊加橙色細線
    rect(slide, Inches(5.75), 0, Inches(0.06), H, t["acc2"])

    # 左側白色主標題
    txt(slide, data.get("title",""), Inches(0.4), Inches(1.5), Inches(5.0), Inches(3.0),
        30, t["title"], bold=True)

    # 橙色分隔線
    rect(slide, Inches(0.4), Inches(4.7), Inches(2.5), Inches(0.06), t["acc2"])

    # 副標題
    txt(slide, data.get("subtitle",""), Inches(0.4), Inches(4.88), Inches(5.0), Inches(0.7),
        13, t["sub"])

    # 右側（紅底）CECLAW
    txt(slide, "CECLAW", Inches(6.1), Inches(1.2), Inches(3.5), Inches(0.65),
        18, RGBColor(0xFF,0xFF,0xFF), bold=True)
    txt(slide, "Enterprise AI Platform", Inches(6.1), Inches(1.8), Inches(3.5), Inches(0.4),
        9, RGBColor(0xFF,0xCC,0xCC))
    txt(slide, "ColdElectric", Inches(6.1), Inches(2.3), Inches(3.5), Inches(0.4),
        10, RGBColor(0xFF,0xAA,0xAA))

    # 右側大數字
    txt(slide, "01", Inches(6.1), H - Inches(1.8), Inches(3.5), Inches(1.4),
        60, RGBColor(0xFF,0xFF,0xFF), bold=True)

    # 底部
    rect(slide, 0, H - Inches(0.1), Inches(5.7), Inches(0.1), t["acc"])
    txt(slide, "數據主權 · 本地推理 · 白名單控制", Inches(0.4), H - Inches(0.42),
        Inches(4.5), Inches(0.32), 8, t["sub"])


def taiwan_content(prs, data, t, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    bg(slide, t["bg"])

    # 頂部紅色條
    rect(slide, 0, 0, W, Inches(0.07), t["acc"])
    # 頂部右側紅色小塊
    rect(slide, W - Inches(1.8), 0, Inches(1.8), Inches(0.75), t["sidebar"])

    # 標題
    rect(slide, Inches(0.4), Inches(0.2), Inches(0.07), Inches(0.6), t["acc2"])
    txt(slide, data.get("title",""), Inches(0.58), Inches(0.16), Inches(7.5), Inches(0.72),
        23, t["title"], bold=True)
    txt(slide, f"{page_num:02d}/{total:02d}", W - Inches(1.65), Inches(0.16),
        Inches(1.5), Inches(0.55), 12, RGBColor(0xFF,0xFF,0xFF), bold=True, align=PP_ALIGN.CENTER)

    # 橫線（橙色）
    rect(slide, Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.025), t["acc2"])

    # 要點
    points = data.get("points", [])
    y_base = Inches(1.15)
    gap = (H - Inches(1.8)) / max(len(points[:4]), 1)

    for i, pt in enumerate(points[:4]):
        y = y_base + i * gap
        num_color = t["acc"] if i % 2 == 0 else t["acc2"]
        # 數字框
        rect(slide, Inches(0.4), y + Inches(0.06), Inches(0.35), Inches(0.35), num_color)
        txt(slide, str(i+1), Inches(0.4), y + Inches(0.06), Inches(0.35), Inches(0.35),
            11, t["num_fg"], bold=True, align=PP_ALIGN.CENTER)
        # 要點文字
        txt(slide, pt, Inches(0.88), y + Inches(0.04), Inches(8.6), gap - Inches(0.12),
            13, t["body"])
        # 分隔虛線（用矩形代替）
        if i < len(points) - 1:
            rect(slide, Inches(0.88), y + gap - Inches(0.06), Inches(8.6), Inches(0.012),
                 t["rule"])

    # 底部
    rect(slide, 0, H - Inches(0.08), W, Inches(0.08), t["acc"])
    txt(slide, "CECLAW · ColdElectric", Inches(0.4), H - Inches(0.4),
        Inches(4), Inches(0.3), 8, t["sub"])


# ══════════════════════════════════════════════════
# ROUTER
# ══════════════════════════════════════════════════
BUILDERS = {
    "dark":   (dark_cover,   dark_content),
    "light":  (light_cover,  light_content),
    "taiwan": (taiwan_cover, taiwan_content),
}


async def generate_pptx(topic: str, pages: int, output_path: str, theme: str = "dark") -> str:
    t = C.get(theme, C["dark"])
    cover_fn, content_fn = BUILDERS.get(theme, BUILDERS["dark"])

    search_results = await search_web(topic)
    slides_data = await llm_generate_slides(topic, pages, search_results)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7)
    total = len(slides_data)

    for i, slide_data in enumerate(slides_data):
        if slide_data.get("type") == "cover" or i == 0:
            cover_fn(prs, slide_data, t)
        else:
            content_fn(prs, slide_data, t, i + 1, total)

    prs.save(output_path)
    return output_path
